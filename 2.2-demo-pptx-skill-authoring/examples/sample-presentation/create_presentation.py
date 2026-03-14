#!/usr/bin/env python3
"""
Sample PPTX presentation created using python-pptx.
Demonstrates the principles in the pptx-presentation skill.

Run: python create_presentation.py
Output: sample-presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ============ SETUP ============

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color palette (PRINCIPLE 1: Color System)
COLOR_PRIMARY = RGBColor(37, 99, 235)      # Blue
COLOR_ACCENT = RGBColor(249, 115, 22)     # Orange
COLOR_TEXT = RGBColor(26, 26, 26)          # Dark gray
COLOR_LIGHT_BG = RGBColor(248, 250, 252)  # Light gray
COLOR_WHITE = RGBColor(255, 255, 255)     # White

# Define spacing (PRINCIPLE 2: Consistent Spacing)
MARGIN = Inches(0.5)
CONTENT_WIDTH = prs.slide_width - (MARGIN * 2)

def set_title_style(shape):
    """Apply consistent title styling"""
    text_frame = shape.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(54)
    paragraph.font.bold = True
    paragraph.font.color.rgb = COLOR_TEXT
    paragraph.alignment = PP_ALIGN.LEFT

def set_body_style(text_frame, size=24):
    """Apply consistent body text styling"""
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = COLOR_TEXT
        paragraph.space_before = Pt(12)
        paragraph.space_after = Pt(12)
        paragraph.line_spacing = 1.4

def add_bullet_slide(title_text, bullets):
    """Add a slide with title and bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE

    # Title
    title_box = slide.shapes.add_textbox(MARGIN, MARGIN, CONTENT_WIDTH, Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    set_title_style(title_box)

    # Bullets
    body_box = slide.shapes.add_textbox(MARGIN, Inches(1.5), CONTENT_WIDTH, Inches(5))
    text_frame = body_box.text_frame
    text_frame.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = bullet
        p.level = 0
        p.font.size = Pt(24)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(8)
        p.space_after = Pt(8)

# ============ SLIDE 1: TITLE SLIDE ============

slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
background = slide1.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = COLOR_PRIMARY

# Title
title_box = slide1.shapes.add_textbox(MARGIN, Inches(2.5), CONTENT_WIDTH, Inches(1.5))
title_frame = title_box.text_frame
title_frame.text = "Skill Authoring Essentials"
title_frame.word_wrap = True
for paragraph in title_frame.paragraphs:
    paragraph.font.size = Pt(60)
    paragraph.font.bold = True
    paragraph.font.color.rgb = COLOR_WHITE
    paragraph.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle_box = slide1.shapes.add_textbox(MARGIN, Inches(4.2), CONTENT_WIDTH, Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Encoding Domain Knowledge with Claude Skills"
for paragraph in subtitle_frame.paragraphs:
    paragraph.font.size = Pt(28)
    paragraph.font.color.rgb = COLOR_LIGHT_BG
    paragraph.alignment = PP_ALIGN.CENTER

# ============ SLIDE 2: WHAT ARE SKILLS? ============

add_bullet_slide(
    "What Are Skills?",
    [
        "Reusable knowledge files that teach Claude best practices",
        "Encode domain expertise in actionable patterns",
        "Guide Claude's behavior toward excellence",
        "Composable: use multiple skills together",
        "Versionable: iterate and improve over time"
    ]
)

# ============ SLIDE 3: KEY PRINCIPLES ============

slide3 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide3.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = COLOR_WHITE

# Title
title_box = slide3.shapes.add_textbox(MARGIN, MARGIN, CONTENT_WIDTH, Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "Skill Components"
set_title_style(title_box)

# 3-column layout
col_width = (CONTENT_WIDTH - Inches(0.3)) / 3
col_height = Inches(5)

components = [
    ("Purpose", "Why does this skill exist?"),
    ("Principles", "What are the core rules?"),
    ("Patterns", "How do you implement it?")
]

for i, (title, desc) in enumerate(components):
    x = MARGIN + i * (col_width + Inches(0.15))

    # Background box
    shape = slide3.shapes.add_shape(
        1,  # Rectangle
        x, Inches(1.5), col_width, col_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_LIGHT_BG
    shape.line.color.rgb = COLOR_PRIMARY
    shape.line.width = Pt(2)

    # Component title
    comp_title = slide3.shapes.add_textbox(x + Inches(0.2), Inches(1.7), col_width - Inches(0.4), Inches(0.5))
    comp_frame = comp_title.text_frame
    comp_frame.text = title
    comp_frame.word_wrap = True
    for p in comp_frame.paragraphs:
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY

    # Component description
    comp_desc = slide3.shapes.add_textbox(x + Inches(0.2), Inches(2.4), col_width - Inches(0.4), Inches(3.5))
    desc_frame = comp_desc.text_frame
    desc_frame.text = desc
    desc_frame.word_wrap = True
    for p in desc_frame.paragraphs:
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXT

# ============ SLIDE 4: ANATOMY OF A SKILL FILE ============

add_bullet_slide(
    "Anatomy of a Skill File",
    [
        "# Skill Name — Clear, descriptive heading",
        "## Purpose — What does this skill solve?",
        "## Principles — 5-7 core guidelines or rules",
        "## Implementation Checklist — How to verify success",
        "## Example Patterns — Code snippets and templates",
        "## When to Apply — Usage scenarios"
    ]
)

# ============ SLIDE 5: SKILL AUTHORING WORKFLOW ============

add_bullet_slide(
    "Creating Your Own Skill",
    [
        "1. Identify the domain (frontend design, API design, testing)",
        "2. Extract principles (what makes good output?)",
        "3. Document patterns (show code examples)",
        "4. List anti-patterns (what mistakes to avoid?)",
        "5. Create before/after examples",
        "6. Test and refine with Claude"
    ]
)

# ============ SLIDE 6: CLOSING ============

slide6 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide6.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = COLOR_PRIMARY

# Main text
main_box = slide6.shapes.add_textbox(MARGIN, Inches(2.5), CONTENT_WIDTH, Inches(2))
main_frame = main_box.text_frame
main_frame.text = "Skills are the future of AI-assisted development"
main_frame.word_wrap = True
for p in main_frame.paragraphs:
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

# Call to action
cta_box = slide6.shapes.add_textbox(MARGIN, Inches(4.5), CONTENT_WIDTH, Inches(1))
cta_frame = cta_box.text_frame
cta_frame.text = "Now create your own skill!"
for p in cta_frame.paragraphs:
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_ACCENT
    p.alignment = PP_ALIGN.CENTER

# ============ SAVE ============

output_path = "sample-presentation.pptx"
prs.save(output_path)
print(f"✓ Presentation created: {output_path}")
print(f"✓ Total slides: {len(prs.slides)}")
print(f"✓ Demonstrates: typography, color system, spacing, composition")
