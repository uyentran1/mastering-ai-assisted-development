from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- Color System ---
PRIMARY = RGBColor(37, 99, 235)       # Brand blue
SECONDARY = RGBColor(245, 158, 11)    # Warm amber accent
DARK_TEXT = RGBColor(26, 26, 26)
BODY_TEXT = RGBColor(71, 85, 105)
WHITE = RGBColor(255, 255, 255)
LIGHT_BG = RGBColor(248, 250, 252)

# --- Helper Functions ---
def style_title(shape, text, size=44):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT
    p.font.name = "Calibri"

def style_body(text_frame, text, size=20):
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = BODY_TEXT
    p.font.name = "Calibri"

def add_accent_bar(slide, left, top, width, height, color=PRIMARY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_bullet_slide(prs, title_text, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    # Accent bar
    add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.5), Inches(1))
    style_title(title_box, title_text, size=40)
    # Bullets
    body_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(8.0), Inches(5.0))
    tf = body_box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(22)
        p.font.color.rgb = BODY_TEXT
        p.font.name = "Calibri"
        p.space_after = Pt(14)
        p.level = 0
    return slide

# --- Create Presentation ---
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ============================================
# SLIDE 1: Title Slide
# ============================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

# Background accent block
bg_shape = slide1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(3.5)
)
bg_shape.fill.solid()
bg_shape.fill.fore_color.rgb = PRIMARY
bg_shape.line.fill.background()

# Title
title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(8.5), Inches(1.5))
tf = title_box.text_frame
tf.clear()
p = tf.paragraphs[0]
p.text = "The Wonderful World of Cats"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = "Calibri"

# Subtitle
sub_box = slide1.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(8.5), Inches(0.8))
tf = sub_box.text_frame
p = tf.paragraphs[0]
p.text = "Why cats have captivated humans for millennia"
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(219, 234, 254)
p.font.name = "Calibri"

# Cat emoji placeholder shape
cat_shape = slide1.shapes.add_shape(
    MSO_SHAPE.OVAL, Inches(3.5), Inches(4.2), Inches(3), Inches(3)
)
cat_shape.fill.solid()
cat_shape.fill.fore_color.rgb = SECONDARY
cat_shape.line.fill.background()
cat_tf = cat_shape.text_frame
cat_tf.word_wrap = True
p = cat_tf.paragraphs[0]
p.text = "🐱"
p.font.size = Pt(72)
p.alignment = PP_ALIGN.CENTER

# ============================================
# SLIDE 2: Why Cats? (Context)
# ============================================
add_bullet_slide(prs, "Why Cats?", [
    "🏠  Over 600 million domestic cats worldwide",
    "📜  Domesticated ~10,000 years ago in the Near East",
    "❤️  Most popular pet in the world by count",
    "🧠  Cats reduce stress and lower blood pressure",
    "🌐  Dominate internet culture and social media",
])

# ============================================
# SLIDE 3: Cat Breeds Overview
# ============================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(slide3, Inches(0), Inches(0), Inches(0.15), Inches(7.5))

title_box = slide3.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.5), Inches(1))
style_title(title_box, "Popular Cat Breeds", size=40)

# Three comparison cards
breeds = [
    ("Persian", "Calm & affectionate", "Long, luxurious coat"),
    ("Siamese", "Vocal & social", "Striking blue eyes"),
    ("Maine Coon", "Gentle giants", "Largest domestic breed"),
]

for i, (name, trait, detail) in enumerate(breeds):
    left = Inches(0.8 + i * 3.1)
    # Card background
    card = slide3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.0), Inches(2.7), Inches(4.5)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_BG
    card.line.color.rgb = RGBColor(226, 232, 240)
    card.line.width = Pt(1)

    # Breed name
    name_box = slide3.shapes.add_textbox(left + Inches(0.2), Inches(2.3), Inches(2.3), Inches(0.8))
    tf = name_box.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.font.name = "Calibri"

    # Trait
    trait_box = slide3.shapes.add_textbox(left + Inches(0.2), Inches(3.3), Inches(2.3), Inches(0.6))
    tf = trait_box.text_frame
    p = tf.paragraphs[0]
    p.text = trait
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_TEXT
    p.font.name = "Calibri"

    # Detail
    detail_box = slide3.shapes.add_textbox(left + Inches(0.2), Inches(4.0), Inches(2.3), Inches(0.6))
    tf = detail_box.text_frame
    p = tf.paragraphs[0]
    p.text = detail
    p.font.size = Pt(16)
    p.font.color.rgb = BODY_TEXT
    p.font.name = "Calibri"

# ============================================
# SLIDE 4: Cat Behavior (Solution/Insight)
# ============================================
add_bullet_slide(prs, "Understanding Cat Behavior", [
    "🐾  Purring signals comfort and self-healing",
    "👀  Slow blinks are a sign of trust and love",
    "📦  Boxes provide security and warmth",
    "🌙  Cats are crepuscular — most active at dawn/dusk",
    "🎯  Hunting instincts drive play behavior",
])

# ============================================
# SLIDE 5: Cats by the Numbers (Evidence)
# ============================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(slide5, Inches(0), Inches(0), Inches(0.15), Inches(7.5))

title_box = slide5.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.5), Inches(1))
style_title(title_box, "Cats by the Numbers", size=40)

stats = [
    ("70%", "of their lives sleeping"),
    ("30 mph", "top running speed"),
    ("230", "bones in their body"),
]

for i, (number, label) in enumerate(stats):
    left = Inches(0.8 + i * 3.1)

    # Number
    num_box = slide5.shapes.add_textbox(left, Inches(2.5), Inches(2.7), Inches(1.5))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    # Label
    lbl_box = slide5.shapes.add_textbox(left, Inches(4.0), Inches(2.7), Inches(1.0))
    tf = lbl_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(20)
    p.font.color.rgb = BODY_TEXT
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

# ============================================
# SLIDE 6: How to Be a Great Cat Owner (CTA)
# ============================================
add_bullet_slide(prs, "How to Be a Great Cat Owner", [
    "🍽️  Provide high-quality nutrition",
    "🏥  Schedule regular vet checkups",
    "🧶  Engage in daily interactive play",
    "🪴  Create vertical spaces for climbing",
    "💧  Always keep fresh water available",
])

# ============================================
# SLIDE 7: Closing Slide
# ============================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])

# Background
bg = slide7.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.5), Inches(10), Inches(3)
)
bg.fill.solid()
bg.fill.fore_color.rgb = PRIMARY
bg.line.fill.background()

# Closing text
close_box = slide7.shapes.add_textbox(Inches(1), Inches(3.0), Inches(8), Inches(1.2))
tf = close_box.text_frame
p = tf.paragraphs[0]
p.text = "Every cat is purrfect."
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

sub_close = slide7.shapes.add_textbox(Inches(1), Inches(4.3), Inches(8), Inches(0.8))
tf = sub_close.text_frame
p = tf.paragraphs[0]
p.text = "Adopt, don't shop.  🐾"
p.font.size = Pt(28)
p.font.color.rgb = RGBColor(219, 234, 254)
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

# --- Save ---
output_path = "cats_presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
