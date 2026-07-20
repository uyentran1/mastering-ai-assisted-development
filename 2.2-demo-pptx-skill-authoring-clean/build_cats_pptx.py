"""Builds cats_presentation.pptx following .claude/skills/pptx-presentation.md."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Color system (single brand color + accent, per skill)
PRIMARY = RGBColor(0x2E, 0x4B, 0x3E)      # deep feline green
ACCENT = RGBColor(0xE8, 0x8D, 0x3A)       # warm amber (cat-eye accent)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
MEDIUM_TEXT = RGBColor(0x47, 0x55, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF7, 0xF5, 0xF0)

FONT = "Calibri"


def style_title(shape, text, size=44, color=DARK_TEXT, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.name = FONT
    p.font.color.rgb = color
    p.alignment = align


def style_bullets(placeholder, bullets, size=22, color=MEDIUM_TEXT):
    tf = placeholder.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(size)
        p.font.name = FONT
        p.font.color.rgb = color
        p.line_spacing = 1.4


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_footer(slide, prs, label):
    box = slide.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(0.4), Inches(4), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(12)
    p.font.name = FONT
    p.font.color.rgb = MEDIUM_TEXT


prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ---------------------------------------------------------------
# Slide 1: Title (Opening — establishes context)
# ---------------------------------------------------------------
slide = blank_slide(prs)
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = PRIMARY
bg.line.fill.background()
bg.shadow.inherit = False

accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.6), Inches(2.2), Inches(0.12))
accent_bar.fill.solid()
accent_bar.fill.fore_color.rgb = ACCENT
accent_bar.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.6), Inches(8.6), Inches(1.8))
style_title(title_box, "The Secret Life of Cats", size=54, color=WHITE)

subtitle_box = slide.shapes.add_textbox(Inches(0.7), Inches(3.7), Inches(8.6), Inches(0.8))
p = subtitle_box.text_frame.paragraphs[0]
p.text = "Biology, behavior, and why they rule the internet"
p.font.size = Pt(24)
p.font.name = FONT
p.font.color.rgb = RGBColor(0xE5, 0xE5, 0xE5)

add_notes(slide, "Welcome the audience. Set up cats as a subject worth taking seriously: biology, behavior, and culture. 1 min.")

# ---------------------------------------------------------------
# Slide 2: Why Cats? (Problem/context — why this matters)
# ---------------------------------------------------------------
slide = blank_slide(prs)
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
style_title(title_box, "Why Cats Deserve a Closer Look")

body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(5))
style_bullets(body_box.text_frame if False else body_box, [
    "600 million+ domestic cats live alongside humans worldwide",
    "Domesticated ~10,000 years ago in the Fertile Crescent",
    "Most popular pet in the world by household count",
    "Widely misunderstood as aloof or low-maintenance",
])
add_footer(slide, prs, "The Secret Life of Cats")
add_notes(slide, "Cats are everywhere but poorly understood. Set up the myth we'll challenge. 1-2 min.")

# ---------------------------------------------------------------
# Slide 3: Anatomy & Senses (Solution/approach — the science)
# ---------------------------------------------------------------
slide = blank_slide(prs)
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
style_title(title_box, "Built for the Hunt")

body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(5))
style_bullets(body_box, [
    "Whiskers sense air currents and gauge tight spaces",
    "Night vision is 6x sharper than a human's",
    "Retractable claws stay sharp for climbing and defense",
    "Flexible spine allows near-instant righting reflex",
    "Ears rotate independently, each up to 180 degrees",
])
add_footer(slide, prs, "The Secret Life of Cats")
add_notes(slide, "Walk through the physical adaptations that make cats efficient predators. 2 min.")

# ---------------------------------------------------------------
# Slide 4: Comparison slide (indoor vs outdoor cats)
# ---------------------------------------------------------------
slide = blank_slide(prs)
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
style_title(title_box, "Indoor vs. Outdoor Living")

col_w = Inches(4.2)
col_h = Inches(4.5)
left_x = Inches(0.5)
right_x = Inches(5.3)
top_y = Inches(1.8)

left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x, top_y, col_w, col_h)
left_box.fill.solid()
left_box.fill.fore_color.rgb = LIGHT_BG
left_box.line.color.rgb = PRIMARY
left_box.line.width = Pt(1.5)
tf = left_box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.3)
p = tf.paragraphs[0]
p.text = "Indoor"
p.font.bold = True
p.font.size = Pt(26)
p.font.name = FONT
p.font.color.rgb = PRIMARY
for line in ["Average lifespan: 12-18 years", "Protected from predators & traffic", "Needs enrichment to stay active"]:
    para = tf.add_paragraph()
    para.text = line
    para.font.size = Pt(18)
    para.font.name = FONT
    para.font.color.rgb = MEDIUM_TEXT
    para.line_spacing = 1.3

right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x, top_y, col_w, col_h)
right_box.fill.solid()
right_box.fill.fore_color.rgb = LIGHT_BG
right_box.line.color.rgb = ACCENT
right_box.line.width = Pt(1.5)
tf = right_box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.3)
p = tf.paragraphs[0]
p.text = "Outdoor"
p.font.bold = True
p.font.size = Pt(26)
p.font.name = FONT
p.font.color.rgb = ACCENT
for line in ["Average lifespan: 2-5 years", "Full range of natural behaviors", "Exposed to disease, traffic, predators"]:
    para = tf.add_paragraph()
    para.text = line
    para.font.size = Pt(18)
    para.font.name = FONT
    para.font.color.rgb = MEDIUM_TEXT
    para.line_spacing = 1.3

add_footer(slide, prs, "The Secret Life of Cats")
add_notes(slide, "Present both lifestyles with equal visual weight, let the lifespan data speak. 1-2 min.")

# ---------------------------------------------------------------
# Slide 5: Data slide (sleep breakdown as a simple bar chart via shapes)
# ---------------------------------------------------------------
slide = blank_slide(prs)
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
style_title(title_box, "A Cat's Day, By the Numbers")

categories = [("Sleeping", 16, PRIMARY), ("Grooming", 4, ACCENT), ("Playing/Hunting", 2, RGBColor(0x8A, 0xA6, 0x99)), ("Eating", 2, RGBColor(0xC9, 0xB4, 0x8F))]
max_hours = 16
chart_left = Inches(1.2)
chart_top = Inches(2.0)
chart_width = Inches(7.5)
bar_height = Inches(0.7)
gap = Inches(0.35)

for i, (label, hours, color) in enumerate(categories):
    y = chart_top + i * (bar_height + gap)
    bar_w = Inches(0.3) + (chart_width - Inches(0.3)) * (hours / max_hours)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, chart_left, y, bar_w, bar_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    label_box = slide.shapes.add_textbox(Inches(0.5), y, Inches(0.7), bar_height)
    p = label_box.text_frame.paragraphs[0]
    p.text = ""
    label_left = slide.shapes.add_textbox(chart_left - Inches(0.05), y - Inches(0.32), Inches(3), Inches(0.3))
    p = label_left.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(16)
    p.font.name = FONT
    p.font.color.rgb = DARK_TEXT
    p.font.bold = True

    value_box = slide.shapes.add_textbox(chart_left + bar_w + Inches(0.1), y, Inches(1.2), bar_height)
    tf2 = value_box.text_frame
    tf2.vertical_anchor = None
    p = tf2.paragraphs[0]
    p.text = f"{hours}h"
    p.font.size = Pt(18)
    p.font.name = FONT
    p.font.color.rgb = MEDIUM_TEXT

add_footer(slide, prs, "The Secret Life of Cats")
add_notes(slide, "Cats sleep 12-16 hours a day, roughly two-thirds of their lives. Ground the fun fact in real data. 1-2 min.")

# ---------------------------------------------------------------
# Slide 6: Closing (Call to Action)
# ---------------------------------------------------------------
slide = blank_slide(prs)
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = PRIMARY
bg.line.fill.background()
bg.shadow.inherit = False

title_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.6), Inches(8.6), Inches(1.2))
style_title(title_box, "Adopt. Observe. Appreciate.", size=44, color=WHITE)

sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(3.7), Inches(8.6), Inches(1))
p = sub_box.text_frame.paragraphs[0]
p.text = "Visit your local shelter and meet a cat looking for a home."
p.font.size = Pt(22)
p.font.name = FONT
p.font.color.rgb = RGBColor(0xE5, 0xE5, 0xE5)

add_notes(slide, "Close with a concrete call to action: shelter visit or adoption. 1 min.")

prs.save("cats_presentation.pptx")
print("Saved cats_presentation.pptx")
